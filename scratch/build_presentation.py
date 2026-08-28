import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    # Load template
    template_path = 'KIT_SPOC_Working_Copy.pptx'
    prs = Presentation(template_path)
    
    # Palette definition
    NAVY_BG = RGBColor(11, 15, 25)
    CARD_BG = RGBColor(20, 28, 48)
    CYAN_ACCENT = RGBColor(0, 229, 255)
    AMBER_ACCENT = RGBColor(255, 107, 0)
    GREEN_ACCENT = RGBColor(0, 255, 157)
    RED_ACCENT = RGBColor(255, 51, 102)
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(200, 210, 225)
    TEXT_DARK = RGBColor(30, 41, 59)
    BORDER_COLOR = RGBColor(0, 180, 220)

    # ----------------------------------------------------
    # SLIDE 1: TITLE PAGE
    # ----------------------------------------------------
    s1 = prs.slides[0]
    
    # Clean default textboxes that overlap
    for shape in list(s1.shapes):
        if shape.name in ['TextBox 9', 'Subtitle 3', 'Title 7', 'Rectangle 24', 'Freeform: Shape 26', 'Picture 4']:
            sp = shape._element
            sp.getparent().remove(sp)
            
    # Add High-tech Title Box
    title_box = s1.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.8))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "SMART INDIA HACKATHON 2026"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = AMBER_ACCENT
    
    p2 = tf1.add_paragraph()
    p2.text = "☣️ CBRS-X : CBRN HAZMAT PROTOCOL & DISASTER SIMULATOR"
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_DARK
    
    p3 = tf1.add_paragraph()
    p3.text = "Enterprise Tactical VR & WebGL Emergency Training Platform with Deterministic SOP Scoring Engine"
    p3.font.size = Pt(13)
    p3.font.italic = True
    p3.font.color.rgb = RGBColor(80, 90, 105)

    # Metadata Card (Left side)
    meta_box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.5), Inches(5.8), Inches(4.3))
    meta_box.fill.solid()
    meta_box.fill.fore_color.rgb = RGBColor(245, 248, 252)
    meta_box.line.color.rgb = RGBColor(0, 120, 215)
    meta_box.line.width = Pt(1.5)
    
    mtf = meta_box.text_frame
    mtf.word_wrap = True
    mtf.margin_left = Inches(0.25)
    mtf.margin_top = Inches(0.25)
    
    mp = mtf.paragraphs[0]
    mp.text = "📋 MISSION METADATA & TEAM DETAILS"
    mp.font.size = Pt(15)
    mp.font.bold = True
    mp.font.color.rgb = RGBColor(0, 80, 160)
    
    meta_items = [
        ("Problem Statement ID:", "SIH260088"),
        ("Problem Statement Title:", "VR / Web-based CBRN Hazmat Response Simulator"),
        ("Theme:", "Disaster Management"),
        ("Category:", "Software / Enterprise Tactical Simulation"),
        ("Institution:", "Kalpataru Institute of Technology (KIT)"),
        ("Team Lead:", "Lohith R C"),
        ("Team Members:", "Monica K S, Chandana M P, Chandana M N, Harshini R B, Pavitra J H")
    ]
    
    for label, val in meta_items:
        p = mtf.add_paragraph()
        p.space_before = Pt(4)
        run1 = p.add_run()
        run1.text = f"{label} "
        run1.font.bold = True
        run1.font.size = Pt(11)
        run1.font.color.rgb = TEXT_DARK
        
        run2 = p.add_run()
        run2.text = val
        run2.font.size = Pt(11)
        run2.font.color.rgb = RGBColor(40, 50, 70)

    # Hero Image (Right side)
    hero_img = "Bay03_Cinematic_Screenshots/01_Wide_Establishing_Shots/03_Bay03_Isometric_Overview.png"
    if os.path.exists(hero_img):
        s1.shapes.add_picture(hero_img, Inches(6.8), Inches(2.5), Inches(5.7), Inches(4.3))

    # ----------------------------------------------------
    # SLIDE 2: PROPOSED SOLUTION & INNOVATION
    # ----------------------------------------------------
    s2 = prs.slides[1]
    
    # Update Title
    for shape in s2.shapes:
        if shape.has_text_frame and "IDEA TITLE" in shape.text:
            shape.text = "PROPOSED SOLUTION & INNOVATION // CBRS-X"
            shape.text_frame.paragraphs[0].font.size = Pt(22)
            shape.text_frame.paragraphs[0].font.bold = True
        elif shape.name == "TextBox 8":
            sp = shape._element
            sp.getparent().remove(sp)
        elif shape.has_text_frame and "Your Team Name" in shape.text:
            shape.text = "CBRS-X Team"

    # 3 Cards Layout
    cards_data = [
        ("🚨 The Critical Problem", [
            "• High physical danger & risk of toxic chemical/radiation exposure.",
            "• Single-use Level-A suits cost ₹50,000+ per drill.",
            "• Classrooms lack realistic tactile and spatial muscle memory.",
            "• Subjective grading without precise SOP compliance timestamps."
        ], RGBColor(255, 240, 243), RGBColor(220, 50, 70)),
        
        ("☣️ The CBRS-X Solution", [
            "• Zero-Risk 3D Simulation of Industrial Storage Bay 03.",
            "• Zero-Install WebGL running smoothly inside any web browser.",
            "• Full 5-Step SOP: Detect ➔ Protect ➔ Contain ➔ Evacuate ➔ Decon.",
            "• Deterministic 100-Point Audit Engine evaluating live telemetry."
        ], RGBColor(240, 249, 255), RGBColor(0, 140, 230)),
        
        ("🌟 Innovation & Uniqueness", [
            "• Dual-Screen Sync: WebGL Trainee Station + Instructor Radar HUD.",
            "• Real-time STOMP WebSockets streaming telemetry at <50ms.",
            "• Tamper-Evident Certification with embedded SHA-256 hash digest.",
            "• Polymorphic Scenario Engine (Chemical, Biological, Radiological)."
        ], RGBColor(240, 253, 244), RGBColor(16, 185, 129))
    ]

    for idx, (title, points, fill_c, line_c) in enumerate(cards_data):
        left = Inches(0.8 + idx * 3.9)
        top = Inches(1.3)
        width = Inches(3.7)
        height = Inches(3.3)
        
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = fill_c
        card.line.color.rgb = line_c
        card.line.width = Pt(1.5)
        
        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Inches(0.18)
        ctf.margin_top = Inches(0.18)
        
        cp = ctf.paragraphs[0]
        cp.text = title
        cp.font.size = Pt(13)
        cp.font.bold = True
        cp.font.color.rgb = line_c
        
        for pt in points:
            p = ctf.add_paragraph()
            p.text = pt
            p.font.size = Pt(10.5)
            p.font.color.rgb = TEXT_DARK
            p.space_before = Pt(3)

    # 5-Stage SOP Process Ribbon across bottom
    stages = [
        ("1. DETECT", "PID Gas Sensor\nPPM Gradients"),
        ("2. PROTECT", "Level-A Hazmat\nSCBA Air Donning"),
        ("3. CONTAIN", "Magnetic Patch\nLeak Torque Seal"),
        ("4. EVACUATE", "Casualty Triage\nWarm Zone Transit"),
        ("5. DECONTAMINATE", "Multi-Stage Arch\nWashdown Verification")
    ]
    
    ribbon_y = Inches(4.75)
    for idx, (st_name, st_desc) in enumerate(stages):
        left = Inches(0.8 + idx * 2.36)
        width = Inches(2.25)
        height = Inches(1.4)
        
        box = s2.shapes.add_shape(MSO_SHAPE.CHEVRON, left, ribbon_y, width, height)
        box.fill.solid()
        colors = [RGBColor(255, 107, 0), RGBColor(230, 140, 0), RGBColor(220, 50, 70), RGBColor(0, 140, 220), RGBColor(16, 185, 129)]
        box.fill.fore_color.rgb = colors[idx]
        box.line.color.rgb = RGBColor(255, 255, 255)
        box.line.width = Pt(1)
        
        btf = box.text_frame
        btf.word_wrap = True
        btf.margin_left = Inches(0.1)
        btf.margin_top = Inches(0.12)
        
        bp1 = btf.paragraphs[0]
        bp1.text = st_name
        bp1.font.size = Pt(11)
        bp1.font.bold = True
        bp1.font.color.rgb = RGBColor(255, 255, 255)
        bp1.alignment = PP_ALIGN.CENTER
        
        bp2 = btf.add_paragraph()
        bp2.text = st_desc
        bp2.font.size = Pt(9)
        bp2.font.color.rgb = RGBColor(255, 255, 255)
        bp2.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # SLIDE 3: TECHNICAL APPROACH & ARCHITECTURE
    # ----------------------------------------------------
    s3 = prs.slides[2]
    for shape in s3.shapes:
        if shape.has_text_frame and "TECHNICAL APPROACH" in shape.text:
            shape.text = "TECHNICAL APPROACH & SYSTEM ARCHITECTURE"
            shape.text_frame.paragraphs[0].font.size = Pt(22)
            shape.text_frame.paragraphs[0].font.bold = True
        elif shape.name == "TextBox 8":
            sp = shape._element
            sp.getparent().remove(sp)
        elif shape.has_text_frame and "Your Team Name" in shape.text:
            shape.text = "CBRS-X Team"

    # Architecture Left Box (Tech Stack Badges)
    tech_box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.3), Inches(4.2), Inches(4.8))
    tech_box.fill.solid()
    tech_box.fill.fore_color.rgb = RGBColor(248, 250, 252)
    tech_box.line.color.rgb = RGBColor(0, 120, 215)
    tech_box.line.width = Pt(1.5)
    
    ttf = tech_box.text_frame
    ttf.word_wrap = True
    ttf.margin_left = Inches(0.2)
    ttf.margin_top = Inches(0.2)
    
    tp = ttf.paragraphs[0]
    tp.text = "🛠️ MULTI-TIER TECHNOLOGY STACK"
    tp.font.size = Pt(13)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(0, 80, 160)
    
    tech_stack = [
        ("Simulation Engine:", "Unity 2022.3 LTS, WebGL / WebAssembly, OpenXR, URP"),
        ("Trainee Web Station:", "React 18, Three.js, Vite 5"),
        ("Command Dashboard:", "React 18 Glassmorphism, Recharts, 2D/3D Radar Map"),
        ("Core Backend API:", "Java 17, Spring Boot 3.2.5, REST Endpoints"),
        ("Live Telemetry:", "STOMP WebSockets (/ws-cbrsx) <50ms latency"),
        ("Database & DevOps:", "PostgreSQL 15, Docker Compose, Nginx Reverse Proxy"),
        ("Security Hardening:", "SHA-256 PDF Digest, LRU Rate Limiting (10k IP)")
    ]
    
    for title, desc in tech_stack:
        p = ttf.add_paragraph()
        p.space_before = Pt(3)
        r1 = p.add_run()
        r1.text = f"{title}\n"
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = RGBColor(15, 23, 42)
        
        r2 = p.add_run()
        r2.text = f"{desc}"
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = RGBColor(70, 80, 95)

    # Architecture Right Diagram (Flow blocks)
    diag_box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.3), Inches(7.3), Inches(4.8))
    diag_box.fill.solid()
    diag_box.fill.fore_color.rgb = RGBColor(240, 245, 255)
    diag_box.line.color.rgb = RGBColor(0, 140, 230)
    diag_box.line.width = Pt(1.5)
    
    dtf = diag_box.text_frame
    dtf.word_wrap = True
    dtf.margin_left = Inches(0.2)
    dtf.margin_top = Inches(0.2)
    
    dp = dtf.paragraphs[0]
    dp.text = "📡 END-TO-END TELEMETRY & SCORING PIPELINE"
    dp.font.size = Pt(13)
    dp.font.bold = True
    dp.font.color.rgb = RGBColor(0, 100, 200)

    # 3 Pipeline Steps
    pipe_steps = [
        ("1. Frontline Ingestion (WebGL Client)", "Trainee executes SOP actions (gas detection, suit donning, containment). CbrsEventLogger transmits ISO-8601 UTC JSON payloads via REST & WebSockets."),
        ("2. Deterministic Audit Core (Spring Boot 3.2.5)", "Evaluates 100-point protocol matrix with millisecond timestamps. Detects safety violations (e.g. entering hot zone without PPE) & calculates real-time deductions."),
        ("3. Tactical Broadcast & Tamper-Proof Credentialing", "Streams live responder coordinates to Instructor Radar HUD at 60Hz. On completion, dynamically generates an official PDF certificate with embedded SHA-256 hash.")
    ]
    
    for title, desc in pipe_steps:
        p = dtf.add_paragraph()
        p.space_before = Pt(6)
        r1 = p.add_run()
        r1.text = f"{title}\n"
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = RGBColor(10, 30, 60)
        
        r2 = p.add_run()
        r2.text = f"{desc}"
        r2.font.size = Pt(10)
        r2.font.color.rgb = RGBColor(60, 70, 85)

    # Verification Badge at bottom of slide 3
    vbox = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.4), Inches(4.9), Inches(6.9), Inches(0.9))
    vbox.fill.solid()
    vbox.fill.fore_color.rgb = RGBColor(16, 185, 129)
    vbox.line.color.rgb = RGBColor(255, 255, 255)
    vtf = vbox.text_frame
    vtf.margin_top = Inches(0.12)
    vp = vtf.paragraphs[0]
    vp.text = "✔ VERIFIED RELIABILITY: 46 / 46 Unit & Integration Tests Passed (JUnit 5 + MockMvc)"
    vp.font.bold = True
    vp.font.size = Pt(11.5)
    vp.font.color.rgb = RGBColor(255, 255, 255)
    vp.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # ----------------------------------------------------
    s4 = prs.slides[3]
    for shape in s4.shapes:
        if shape.has_text_frame and "FEASIBILITY AND VIABILITY" in shape.text:
            shape.text = "FEASIBILITY, RISK ANALYSIS & MITIGATION"
            shape.text_frame.paragraphs[0].font.size = Pt(22)
            shape.text_frame.paragraphs[0].font.bold = True
        elif shape.name == "TextBox 8":
            sp = shape._element
            sp.getparent().remove(sp)
        elif shape.has_text_frame and "Your Team Name" in shape.text:
            shape.text = "CBRS-X Team"

    # Risk Table Card (Left)
    rtable_box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.3), Inches(6.5), Inches(4.8))
    rtable_box.fill.solid()
    rtable_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    rtable_box.line.color.rgb = RGBColor(220, 50, 70)
    rtable_box.line.width = Pt(1.5)
    
    rtf = rtable_box.text_frame
    rtf.word_wrap = True
    rtf.margin_left = Inches(0.2)
    rtf.margin_top = Inches(0.2)
    
    rp = rtf.paragraphs[0]
    rp.text = "⚠️ RISK ANALYSIS & ENGINEERING MITIGATION"
    rp.font.size = Pt(13)
    rp.font.bold = True
    rp.font.color.rgb = RGBColor(200, 30, 50)

    risks = [
        ("Heavy 3D WebGL Bundle Sizes", "Brotli compression + WebAssembly reduces initial payload by 65%. Service Worker caching enables instant repeat drills."),
        ("No Internet in Remote Bases", "Offline 'Field Box' Deployment: Entire stack packaged via Docker Compose running on a local offline Wi-Fi access point."),
        ("Credential Forgery & Tampering", "Embedded SHA-256 cryptographic digest in PDF certificates allows instant offline hash validation."),
        ("Cybersecurity & Ingress Threats", "Bounded LRU Sliding-Window Rate Limiter (10,000 capacity) + DDE / CSV injection sanitization filters.")
    ]

    for title, desc in risks:
        p = rtf.add_paragraph()
        p.space_before = Pt(4)
        r1 = p.add_run()
        r1.text = f"• {title}: "
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = RGBColor(15, 23, 42)
        
        r2 = p.add_run()
        r2.text = f"{desc}"
        r2.font.size = Pt(10)
        r2.font.color.rgb = RGBColor(70, 80, 95)

    # Right Card: Hardware & Deployment Feasibility
    feat_box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.3), Inches(5.0), Inches(4.8))
    feat_box.fill.solid()
    feat_box.fill.fore_color.rgb = RGBColor(245, 248, 255)
    feat_box.line.color.rgb = RGBColor(0, 120, 215)
    feat_box.line.width = Pt(1.5)
    
    ftf = feat_box.text_frame
    ftf.word_wrap = True
    ftf.margin_left = Inches(0.2)
    ftf.margin_top = Inches(0.2)
    
    fp = ftf.paragraphs[0]
    fp.text = "📦 PORTABLE 'FIELD BOX' DEPLOYMENT"
    fp.font.size = Pt(13)
    fp.font.bold = True
    fp.font.color.rgb = RGBColor(0, 80, 180)

    fpoints = [
        ("Zero Hardware Dependency:", "Runs inside standard Google Chrome / Microsoft Edge on existing budget laptops without dedicated GPUs."),
        ("Single-Command Field Setup:", "`docker-compose up` launches Nginx, Spring Boot backend, and PostgreSQL in under 15 seconds."),
        ("Concurrent Trainee Support:", "Handles up to 50 simultaneous browser simulation streams per local base server without frame drops."),
        ("Multi-Scenario Polymorphism:", "Hot-swappable JSON state machines for Chemical, Biological, and Radiological hazards.")
    ]

    for title, desc in fpoints:
        p = ftf.add_paragraph()
        p.space_before = Pt(4)
        r1 = p.add_run()
        r1.text = f"✔ {title}\n"
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = RGBColor(10, 40, 80)
        
        r2 = p.add_run()
        r2.text = f"{desc}"
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = RGBColor(60, 70, 85)

    # ----------------------------------------------------
    # SLIDE 5: IMPACT AND BENEFITS
    # ----------------------------------------------------
    s5 = prs.slides[4]
    for shape in s5.shapes:
        if shape.has_text_frame and "IMPACT AND BENEFITS" in shape.text:
            shape.text = "REAL-WORLD IMPACT, ROI & SCALABILITY"
            shape.text_frame.paragraphs[0].font.size = Pt(22)
            shape.text_frame.paragraphs[0].font.bold = True
        elif shape.name == "TextBox 8":
            sp = shape._element
            sp.getparent().remove(sp)
        elif shape.has_text_frame and "Your Team Name" in shape.text:
            shape.text = "CBRS-X Team"

    # 3 High-Impact Metric Cards
    metrics_data = [
        ("🚨 100% RISK ELIMINATION", "Zero physical injury, chemical burns, or toxic vapor inhalation during emergency training.", RGBColor(255, 235, 238), RGBColor(220, 50, 70)),
        ("💰 90%+ COST REDUCTION", "Saves ₹50k–₹1L per responder per drill by eliminating single-use Level-A suit & reagent waste.", RGBColor(235, 255, 245), RGBColor(16, 185, 129)),
        ("📈 1,000+ DRILLS / DAY", "Infinite virtual repetitions with instantaneous objective grading and replay analysis.", RGBColor(235, 248, 255), RGBColor(0, 140, 230))
    ]

    for idx, (title, desc, fill_c, line_c) in enumerate(metrics_data):
        left = Inches(0.8 + idx * 3.9)
        card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.3), Inches(3.7), Inches(1.6))
        card.fill.solid()
        card.fill.fore_color.rgb = fill_c
        card.line.color.rgb = line_c
        card.line.width = Pt(1.5)
        
        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Inches(0.15)
        ctf.margin_top = Inches(0.15)
        
        cp = ctf.paragraphs[0]
        cp.text = title
        cp.font.size = Pt(13)
        cp.font.bold = True
        cp.font.color.rgb = line_c
        
        cp2 = ctf.add_paragraph()
        cp2.text = desc
        cp2.font.size = Pt(10)
        cp2.font.color.rgb = TEXT_DARK
        cp2.space_before = Pt(3)

    # Lower Split: Target Institutions (Left) + Containment Action Image (Right)
    target_box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.1), Inches(6.0), Inches(3.0))
    target_box.fill.solid()
    target_box.fill.fore_color.rgb = RGBColor(248, 250, 252)
    target_box.line.color.rgb = RGBColor(0, 120, 215)
    target_box.line.width = Pt(1.5)
    
    ttf = target_box.text_frame
    ttf.word_wrap = True
    ttf.margin_left = Inches(0.2)
    ttf.margin_top = Inches(0.2)
    
    tp = ttf.paragraphs[0]
    tp.text = "🎯 TARGET AUDIENCE & INSTITUTIONAL REACH"
    tp.font.size = Pt(13)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(0, 80, 160)

    target_items = [
        ("National & State Disaster Units:", "NDRF (National Disaster Response Force), SDRF battalions across India."),
        ("Industrial & Energy Sectors:", "Petrochemical depots, refineries, and chemical plants (IOCL, ONGC, GAIL, NTPC)."),
        ("Civil Defense & Paramilitary:", "CRPF, BSF, CISF HAZMAT squads, and Municipal Fire Rescue Services."),
        ("Environmental Sustainability:", "Zero toxic runoff or plastic suit disposal into municipal landfills.")
    ]

    for title, desc in target_items:
        p = ttf.add_paragraph()
        p.space_before = Pt(3)
        r1 = p.add_run()
        r1.text = f"• {title} "
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = RGBColor(15, 23, 42)
        
        r2 = p.add_run()
        r2.text = f"{desc}"
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = RGBColor(60, 70, 85)

    # Action Screenshot (Right)
    action_img = "Bay03_Cinematic_Screenshots/04_FirstPerson_Perspective/04_Trainee_FirstPerson_Containment_Action.png"
    if os.path.exists(action_img):
        s5.shapes.add_picture(action_img, Inches(7.0), Inches(3.1), Inches(5.5), Inches(3.0))

    # ----------------------------------------------------
    # SLIDE 6: RESEARCH AND REFERENCES
    # ----------------------------------------------------
    s6 = prs.slides[5]
    for shape in s6.shapes:
        if shape.has_text_frame and "RESEARCH" in shape.text:
            shape.text = "RESEARCH BASE, COMPLIANCE & REFERENCES"
            shape.text_frame.paragraphs[0].font.size = Pt(22)
            shape.text_frame.paragraphs[0].font.bold = True
        elif shape.name == "TextBox 8":
            sp = shape._element
            sp.getparent().remove(sp)
        elif shape.has_text_frame and "Your Team Name" in shape.text:
            shape.text = "CBRS-X Team"

    ref_box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.7), Inches(4.8))
    ref_box.fill.solid()
    ref_box.fill.fore_color.rgb = RGBColor(248, 250, 252)
    ref_box.line.color.rgb = RGBColor(0, 120, 215)
    ref_box.line.width = Pt(1.5)
    
    rtf = ref_box.text_frame
    rtf.word_wrap = True
    rtf.margin_left = Inches(0.25)
    rtf.margin_top = Inches(0.25)
    
    rp = rtf.paragraphs[0]
    rp.text = "📑 OFFICIAL REGULATORY STANDARDS & SCIENTIFIC CITATIONS"
    rp.font.size = Pt(14)
    rp.font.bold = True
    rp.font.color.rgb = RGBColor(0, 80, 160)

    references = [
        ("1. NDRF Standard Operating Procedures (SOP):", "Tactical First-Responder Protocol for Chemical, Biological, Radiological and Nuclear (CBRN) Disasters (Govt. of India)."),
        ("2. National Disaster Management Authority (NDMA):", "Management of Chemical (Terrorism) Disasters & Industrial Hazardous Material Response Guidelines (MHA)."),
        ("3. OSHA & NIOSH Hazmat Exposure Limits:", "PID Sensor Calibration Protocols, Permissible Exposure Limits (PEL) & Immediately Dangerous to Life or Health (IDLH) Caps."),
        ("4. Khronos Group & W3C Web Standards:", "OpenXR 1.0 Cross-Platform VR Standard & High-Performance WebGL 2.0 / WebAssembly 3D Pipeline."),
        ("5. Enterprise Security Standards:", "Spring Security 6, STOMP RFC WebSocket Specifications & SHA-256 Cryptographic Hash Digests.")
    ]

    for title, desc in references:
        p = rtf.add_paragraph()
        p.space_before = Pt(7)
        r1 = p.add_run()
        r1.text = f"{title}\n"
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = RGBColor(10, 30, 60)
        
        r2 = p.add_run()
        r2.text = f"{desc}"
        r2.font.size = Pt(10)
        r2.font.color.rgb = RGBColor(60, 70, 85)

    # ----------------------------------------------------
    # REMOVE SLIDE 7 (INSTRUCTION SLIDE AS PER SIH MANDATE)
    # ----------------------------------------------------
    if len(prs.slides) > 6:
        # Delete 7th slide
        rId = prs.slides._sldIdLst[6].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[6]

    # Save output
    output_filename = "CBRS_X_SIH_2026_Idea_Presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation successfully saved to: {output_filename}")

    # Also try to overwrite original if possible
    try:
        prs.save("KIT SPOC 2026-IDEA-Presentation-Format.pptx")
        print("Updated KIT SPOC 2026-IDEA-Presentation-Format.pptx successfully.")
    except Exception as e:
        print("Original template locked by PowerPoint. Generated clean copy instead:", output_filename)

if __name__ == "__main__":
    create_deck()
