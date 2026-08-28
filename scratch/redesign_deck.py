import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_impressive_deck():
    template_path = 'KIT_SPOC_Working_Copy.pptx'
    prs = Presentation(template_path)
    
    # ---------------------------------------------------------
    # COLOR DESIGN TOKENS
    # ---------------------------------------------------------
    DARK_NAVY   = RGBColor(10, 16, 30)      # #0A101E
    SLATE_CARD  = RGBColor(246, 249, 253)   # Crisp clean card bg
    CARD_BORDER = RGBColor(0, 140, 230)     # Primary tech blue
    CYAN_NEON   = RGBColor(0, 180, 220)     # Tech cyan
    AMBER_ALERT = RGBColor(230, 95, 0)      # High-vis Hazmat orange
    GREEN_VALID = RGBColor(16, 160, 100)    # Success / verified green
    RED_HOT     = RGBColor(220, 38, 38)     # Hot zone danger red
    PURPLE_CERT = RGBColor(120, 50, 200)    # Cryptographic cert purple
    TEXT_MAIN   = RGBColor(15, 23, 42)      # Deep slate for primary text
    TEXT_MUTED  = RGBColor(71, 85, 105)     # Secondary slate
    WHITE       = RGBColor(255, 255, 255)

    def add_header_strip(slide, title_text, category_tag="SIH 2026 // DISASTER MANAGEMENT"):
        # Top accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.4), Inches(11.73), Inches(0.7))
        bar.fill.solid()
        bar.fill.fore_color.rgb = DARK_NAVY
        bar.line.color.rgb = CARD_BORDER
        bar.line.width = Pt(1.2)
        
        btf = bar.text_frame
        btf.margin_left = Inches(0.2)
        btf.margin_top = Inches(0.08)
        
        bp = btf.paragraphs[0]
        r_tag = bp.add_run()
        r_tag.text = f"[{category_tag}]  "
        r_tag.font.size = Pt(9.5)
        r_tag.font.bold = True
        r_tag.font.color.rgb = AMBER_ALERT
        
        r_title = bp.add_run()
        r_title.text = title_text
        r_title.font.size = Pt(14)
        r_title.font.bold = True
        r_title.font.color.rgb = WHITE

    # =========================================================
    # SLIDE 1: TITLE PAGE (COMMAND INGRESS & HERO TWIN)
    # =========================================================
    s1 = prs.slides[0]
    for shape in list(s1.shapes):
        if shape.name in ['TextBox 9', 'Subtitle 3', 'Title 7', 'Rectangle 24', 'Freeform: Shape 26', 'Picture 4']:
            sp = shape._element
            sp.getparent().remove(sp)

    # Hero Banner
    banner = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.45), Inches(11.73), Inches(1.85))
    banner.fill.solid()
    banner.fill.fore_color.rgb = DARK_NAVY
    banner.line.color.rgb = CYAN_NEON
    banner.line.width = Pt(1.5)
    
    btf = banner.text_frame
    btf.margin_left = Inches(0.3)
    btf.margin_top = Inches(0.18)
    
    p1 = btf.paragraphs[0]
    p1.text = "SMART INDIA HACKATHON 2026 // PS: SIH260088"
    p1.font.size = Pt(11)
    p1.font.bold = True
    p1.font.color.rgb = AMBER_ALERT
    
    p2 = btf.add_paragraph()
    p2.text = "☣️ CBRS-X : CBRN HAZMAT PROTOCOL & RESPONSE SIMULATOR"
    p2.font.size = Pt(21)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(3)
    
    p3 = btf.add_paragraph()
    p3.text = "Enterprise Tactical VR & WebGL Chemical, Biological & Radiological Emergency Training with Deterministic Audit Engine"
    p3.font.size = Pt(11)
    p3.font.color.rgb = RGBColor(180, 205, 235)
    p3.space_before = Pt(2)

    # Left: Metadata Card with Clean Highlight Rows
    meta_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.45), Inches(5.6), Inches(4.35))
    meta_card.fill.solid()
    meta_card.fill.fore_color.rgb = SLATE_CARD
    meta_card.line.color.rgb = CARD_BORDER
    meta_card.line.width = Pt(1.5)
    
    mtf = meta_card.text_frame
    mtf.margin_left = Inches(0.25)
    mtf.margin_top = Inches(0.2)
    
    mp = mtf.paragraphs[0]
    mp.text = "📋 MISSION SPECIFICATIONS & TEAM"
    mp.font.size = Pt(13)
    mp.font.bold = True
    mp.font.color.rgb = DARK_NAVY
    
    rows = [
        ("Problem Statement ID", "SIH260088"),
        ("Problem Title", "Virtual Reality & WebGL CBRN Emergency Simulator"),
        ("Ministry / Domain", "Disaster Management (NDRF / MHA)"),
        ("PS Category", "Software & Enterprise Telemetry"),
        ("Institution", "Kalpataru Institute of Technology (KIT)"),
        ("Team Lead", "Lohith R C (Unity / Backend / Telemetry)"),
        ("Core Task Force", "Monica K S, Chandana M P, Chandana M N, Harshini R B, Pavitra J H")
    ]
    
    for label, val in rows:
        p = mtf.add_paragraph()
        p.space_before = Pt(4)
        r1 = p.add_run()
        r1.text = f"• {label}: "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = CARD_BORDER
        
        r2 = p.add_run()
        r2.text = val
        r2.font.size = Pt(10)
        r2.font.color.rgb = TEXT_MAIN

    # Right: Hero 3D Picture with HUD Frame
    hero_img = "Bay03_Cinematic_Screenshots/01_Wide_Establishing_Shots/03_Bay03_Isometric_Overview.png"
    if os.path.exists(hero_img):
        s1.shapes.add_picture(hero_img, Inches(6.6), Inches(2.45), Inches(5.93), Inches(4.35))
        
        # HUD Tag overlay on image
        hud_tag = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(6.35), Inches(5.53), Inches(0.35))
        hud_tag.fill.solid()
        hud_tag.fill.fore_color.rgb = DARK_NAVY
        hud_tag.line.fill.background()
        htf = hud_tag.text_frame
        htf.margin_top = Inches(0.04)
        hp = htf.paragraphs[0]
        hp.text = "📍 INDUSTRIAL STORAGE BAY 03 // 3D TACTICAL DIGITAL TWIN"
        hp.font.size = Pt(9.5)
        hp.font.bold = True
        hp.font.color.rgb = GREEN_VALID
        hp.alignment = PP_ALIGN.CENTER

    # =========================================================
    # SLIDE 2: PROPOSED SOLUTION & 5-STAGE SOP
    # =========================================================
    s2 = prs.slides[1]
    for shape in list(s2.shapes):
        if shape.name in ["TextBox 8", "Title 1"]:
            sp = shape._element
            sp.getparent().remove(sp)
        elif shape.has_text_frame and "Your Team Name" in shape.text:
            shape.text = "CBRS-X Team"

    add_header_strip(s2, "PROPOSED SOLUTION & INNOVATION // ZERO-RISK CBRN AUDITING", "MISSION DOMAIN // SIH260088")

    # 3 Distinct Glassmorphic Pillars
    col_w = Inches(3.75)
    gap = Inches(0.24)
    top_y = Inches(1.25)
    card_h = Inches(3.45)
    
    pills = [
        ("🚨 The Real-World Crisis", [
            ("Lethal Risk:", "Live hazardous chemical/radiation exposure causes fatal injury."),
            ("Prohibitive Cost:", "Single-use Level-A suits cost ₹50,000+ per drill."),
            ("Zero Muscle Memory:", "Passive classroom lectures fail to build real reflex."),
            ("Subjective Bias:", "Human evaluations miss exact millisecond protocol errors.")
        ], RED_HOT, RGBColor(255, 245, 245)),
        
        ("☣️ The CBRS-X Solution", [
            ("Zero-Risk 3D Simulation:", "High-fidelity virtual recreation of Storage Bay 03."),
            ("Zero-Install WebGL:", "Runs natively in Chrome/Edge on standard office laptops."),
            ("Full Tactical SOP:", "Enforces official 5-stage NDRF emergency cycle."),
            ("Deterministic Grading:", "100-point audit engine eliminates subjective bias.")
        ], CARD_BORDER, RGBColor(245, 250, 255)),
        
        ("🌟 Core Innovation & Edge", [
            ("Dual-Screen Ecosystem:", "Trainee WebGL view synced live with Instructor Radar."),
            ("Sub-50ms Telemetry:", "STOMP WebSockets stream responder data at 60Hz."),
            ("SHA-256 Certification:", "Instant tamper-evident official PDF credentialing."),
            ("Multi-Hazard Engine:", "Hot-swappable Chemical, Bio & Radiological logic.")
        ], GREEN_VALID, RGBColor(245, 254, 248))
    ]

    for idx, (head, items, border_c, bg_c) in enumerate(pills):
        x = Inches(0.8 + idx * (col_w.inches + gap.inches))
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top_y, col_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_c
        card.line.color.rgb = border_c
        card.line.width = Pt(1.5)
        
        ctf = card.text_frame
        ctf.margin_left = Inches(0.18)
        ctf.margin_top = Inches(0.15)
        
        cp = ctf.paragraphs[0]
        cp.text = head
        cp.font.size = Pt(12)
        cp.font.bold = True
        cp.font.color.rgb = border_c
        
        for k, v in items:
            p = ctf.add_paragraph()
            p.space_before = Pt(4)
            r1 = p.add_run()
            r1.text = f"• {k} "
            r1.font.bold = True
            r1.font.size = Pt(9.5)
            r1.font.color.rgb = TEXT_MAIN
            
            r2 = p.add_run()
            r2.text = v
            r2.font.size = Pt(9)
            r2.font.color.rgb = TEXT_MUTED

    # Bottom: 5-Stage SOP Connected Milestone Ribbon
    stages = [
        ("STAGE 1", "DETECT", "PID Gas Sampling\nPPM Thresholds", AMBER_ALERT),
        ("STAGE 2", "PROTECT", "Level-A Hazmat\nSCBA Air Donning", RGBColor(220, 120, 0)),
        ("STAGE 3", "CONTAIN", "Magnetic Patch\nTorque Leak Seal", RED_HOT),
        ("STAGE 4", "EVACUATE", "Casualty Triage\nWarm Zone Transit", CARD_BORDER),
        ("STAGE 5", "DECONTAMINATE", "Multi-Stage Arch\nWashdown Check", GREEN_VALID)
    ]
    
    sop_y = Inches(4.85)
    st_w = Inches(2.26)
    st_gap = Inches(0.11)
    
    for idx, (num, title, sub, color) in enumerate(stages):
        x = Inches(0.8 + idx * (st_w.inches + st_gap.inches))
        
        box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, sop_y, st_w, Inches(1.35))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = WHITE
        box.line.width = Pt(1)
        
        btf = box.text_frame
        btf.margin_top = Inches(0.08)
        btf.margin_left = Inches(0.08)
        
        bp1 = btf.paragraphs[0]
        bp1.text = f"[{num}]"
        bp1.font.size = Pt(8.5)
        bp1.font.bold = True
        bp1.font.color.rgb = RGBColor(240, 240, 240)
        bp1.alignment = PP_ALIGN.CENTER
        
        bp2 = btf.add_paragraph()
        bp2.text = title
        bp2.font.size = Pt(11)
        bp2.font.bold = True
        bp2.font.color.rgb = WHITE
        bp2.alignment = PP_ALIGN.CENTER
        
        bp3 = btf.add_paragraph()
        bp3.text = sub
        bp3.font.size = Pt(8.5)
        bp3.font.color.rgb = RGBColor(245, 245, 245)
        bp3.alignment = PP_ALIGN.CENTER

    # =========================================================
    # SLIDE 3: TECHNICAL APPROACH & DATA FLOW TOPOLOGY
    # =========================================================
    s3 = prs.slides[2]
    for shape in list(s3.shapes):
        if shape.name in ["TextBox 8", "Title 1"]:
            sp = shape._element
            sp.getparent().remove(sp)
        elif shape.has_text_frame and "Your Team Name" in shape.text:
            shape.text = "CBRS-X Team"

    add_header_strip(s3, "TECHNICAL ARCHITECTURE & DATA FLOW // MULTI-TIER ENGINE", "FULL-STACK SPECS")

    # Left: Multi-Tier Tech Stack Box
    tech_card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.25), Inches(4.3), Inches(4.95))
    tech_card.fill.solid()
    tech_card.fill.fore_color.rgb = SLATE_CARD
    tech_card.line.color.rgb = CARD_BORDER
    tech_card.line.width = Pt(1.5)
    
    ttf = tech_card.text_frame
    ttf.margin_left = Inches(0.2)
    ttf.margin_top = Inches(0.18)
    
    tp = ttf.paragraphs[0]
    tp.text = "🛠️ MULTI-TIER TECHNOLOGY STACK"
    tp.font.size = Pt(12.5)
    tp.font.bold = True
    tp.font.color.rgb = DARK_NAVY

    techs = [
        ("Frontline Simulation:", "Unity 2022.3 LTS, WebGL / WASM, OpenXR, Universal Render Pipeline (URP)"),
        ("Trainee Web Station:", "React 18, Three.js 3D viewport, Vite 5"),
        ("Instructor Command:", "React 18 Glassmorphism, Recharts, 2D/3D Spatial Radar Map"),
        ("Deterministic Core:", "Java 17, Spring Boot 3.2.5, REST API Controllers"),
        ("Real-Time Telemetry:", "STOMP WebSockets (/ws-cbrsx) with <50ms dispatch latency"),
        ("Data Persistence:", "PostgreSQL 15 / Supabase relational audit trails"),
        ("Security Hardening:", "LRU Sliding Rate Limiter (10k IP), SHA-256 PDF engine")
    ]
    
    for t_name, t_desc in techs:
        p = ttf.add_paragraph()
        p.space_before = Pt(3.5)
        r1 = p.add_run()
        r1.text = f"• {t_name} "
        r1.font.bold = True
        r1.font.size = Pt(9.5)
        r1.font.color.rgb = CARD_BORDER
        
        r2 = p.add_run()
        r2.text = t_desc
        r2.font.size = Pt(8.8)
        r2.font.color.rgb = TEXT_MUTED

    # Right Top: Pipeline Diagram Cards
    right_x = Inches(5.3)
    right_w = Inches(7.23)
    
    flow_card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, Inches(1.25), right_w, Inches(3.85))
    flow_card.fill.solid()
    flow_card.fill.fore_color.rgb = RGBColor(240, 247, 255)
    flow_card.line.color.rgb = CARD_BORDER
    flow_card.line.width = Pt(1.5)
    
    ftf = flow_card.text_frame
    ftf.margin_left = Inches(0.2)
    ftf.margin_top = Inches(0.18)
    
    fp = ftf.paragraphs[0]
    fp.text = "📡 END-TO-END TELEMETRY & DETERMINISTIC PIPELINE"
    fp.font.size = Pt(12.5)
    fp.font.bold = True
    fp.font.color.rgb = DARK_NAVY

    pipe_nodes = [
        ("1. Frontline Event Emission (WebGL / Unity C#):", "As trainee navigates Bay 03, CbrsEventLogger transmits ISO-8601 UTC timestamped JSON telemetry (gas detection, suit donning, leak containment, decon exit) over HTTP/WS."),
        ("2. Real-Time Ingress & State Evaluation (Spring Boot 3.2.5):", "Ingress security filters validate payloads; ScoringService deterministically computes deductions for PPE breaches, PPM exposure, and containment speed with 0% subjective bias."),
        ("3. Tactical Broadcast & Tamper-Evident Certification:", "Broadcasts live responder radar coordinates to Instructor Command at 60Hz. At mission end, dynamically generates an official PDF certificate with embedded SHA-256 digest.")
    ]

    for title, desc in pipe_nodes:
        p = ftf.add_paragraph()
        p.space_before = Pt(5)
        r1 = p.add_run()
        r1.text = f"{title}\n"
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = DARK_NAVY
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(9)
        r2.font.color.rgb = TEXT_MUTED

    # Right Bottom: Verified Reliability Banner
    test_bar = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, Inches(5.25), right_w, Inches(0.95))
    test_bar.fill.solid()
    test_bar.fill.fore_color.rgb = GREEN_VALID
    test_bar.line.color.rgb = WHITE
    test_bar.line.width = Pt(1)
    
    tbtf = test_bar.text_frame
    tbtf.margin_top = Inches(0.12)
    tbp1 = tbtf.paragraphs[0]
    tbp1.text = "🏆 VERIFIED PRODUCTION RELIABILITY"
    tbp1.font.bold = True
    tbp1.font.size = Pt(11)
    tbp1.font.color.rgb = WHITE
    tbp1.alignment = PP_ALIGN.CENTER
    
    tbp2 = tbtf.add_paragraph()
    tbp2.text = "✔ 46 / 46 Unit & Integration Tests Passed (JUnit 5 + MockMvc)  |  ✔ Sub-50ms WebSocket Broadcast  |  ✔ 60 FPS WebGL Engine"
    tbp2.font.size = Pt(9)
    tbp2.font.color.rgb = RGBColor(240, 255, 245)
    tbp2.alignment = PP_ALIGN.CENTER

    # =========================================================
    # SLIDE 4: FEASIBILITY & OFFLINE FIELD DEPLOYMENT
    # =========================================================
    s4 = prs.slides[3]
    for shape in list(s4.shapes):
        if shape.name in ["TextBox 8", "Title 1"]:
            sp = shape._element
            sp.getparent().remove(sp)
        elif shape.has_text_frame and "Your Team Name" in shape.text:
            shape.text = "CBRS-X Team"

    add_header_strip(s4, "FEASIBILITY, RISK MITIGATION & FIELD DEPLOYMENT", "TACTICAL RESILIENCE")

    # Left: Risk Analysis & Mitigation Matrix
    risk_card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.25), Inches(6.2), Inches(4.95))
    risk_card.fill.solid()
    risk_card.fill.fore_color.rgb = SLATE_CARD
    risk_card.line.color.rgb = RED_HOT
    risk_card.line.width = Pt(1.5)
    
    rtf = risk_card.text_frame
    rtf.margin_left = Inches(0.2)
    rtf.margin_top = Inches(0.18)
    
    rp = rtf.paragraphs[0]
    rp.text = "⚠️ RISK ANALYSIS & SYSTEMIC MITIGATION"
    rp.font.size = Pt(12.5)
    rp.font.bold = True
    rp.font.color.rgb = RED_HOT

    risk_items = [
        ("Heavy 3D WebGL Bundle Size", "Brotli compression + WebAssembly reduces initial load by 65%. Service Worker caching enables instant repeat launches."),
        ("No Internet at Field Outposts", "Offline 'Field Box' Deployment: Entire ecosystem packaged in Docker Compose running on local offline Wi-Fi."),
        ("Certificate Forgery / Fraud", "Embedded SHA-256 cryptographic digest stamped on all generated PDF certificates allows instant hash verification."),
        ("Cybersecurity & Ingress Risks", "Bounded LRU Sliding-Window Rate Limiting (10k IP pool) + DDE / CSV injection sanitization guards.")
    ]

    for title, desc in risk_items:
        p = rtf.add_paragraph()
        p.space_before = Pt(5)
        r1 = p.add_run()
        r1.text = f"• {title}\n"
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = DARK_NAVY
        
        r2 = p.add_run()
        r2.text = f"  ➔ Mitigation: {desc}"
        r2.font.size = Pt(9)
        r2.font.color.rgb = TEXT_MUTED

    # Right: "Field Box" Offline Architecture Box
    field_card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.25), Inches(5.33), Inches(4.95))
    field_card.fill.solid()
    field_card.fill.fore_color.rgb = RGBColor(245, 248, 255)
    field_card.line.color.rgb = CARD_BORDER
    field_card.line.width = Pt(1.5)
    
    ftf = field_card.text_frame
    ftf.margin_left = Inches(0.2)
    ftf.margin_top = Inches(0.18)
    
    fp = ftf.paragraphs[0]
    fp.text = "📦 PORTABLE 'FIELD BOX' DEPLOYMENT"
    fp.font.size = Pt(12.5)
    fp.font.bold = True
    fp.font.color.rgb = DARK_NAVY

    field_specs = [
        ("Zero Hardware Dependency", "Runs inside standard Google Chrome / Microsoft Edge on budget laptops without dedicated gaming GPUs."),
        ("Single-Command Field Setup", "`docker-compose up` spins up Nginx, Spring Boot backend, and PostgreSQL in under 15 seconds."),
        ("Local Hotspot Connectivity", "Responders connect laptops/tablets to a local router with zero internet connection required."),
        ("High Concurrency Capacity", "Handles up to 50 concurrent simulation streams per local base server without frame rate degradation.")
    ]

    for title, desc in field_specs:
        p = ftf.add_paragraph()
        p.space_before = Pt(5)
        r1 = p.add_run()
        r1.text = f"✔ {title}\n"
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = CARD_BORDER
        
        r2 = p.add_run()
        r2.text = f"  {desc}"
        r2.font.size = Pt(9)
        r2.font.color.rgb = TEXT_MUTED

    # =========================================================
    # SLIDE 5: IMPACT, ROI & SCALABILITY
    # =========================================================
    s5 = prs.slides[4]
    for shape in list(s5.shapes):
        if shape.name in ["TextBox 8", "Title 1"]:
            sp = shape._element
            sp.getparent().remove(sp)
        elif shape.has_text_frame and "Your Team Name" in shape.text:
            shape.text = "CBRS-X Team"

    add_header_strip(s5, "REAL-WORLD IMPACT, ROI & NATIONWIDE SCALABILITY", "STRATEGIC OUTCOMES")

    # Top: 3 High-Impact KPI Counter Cards
    kpis = [
        ("100%", "RISK ELIMINATION", "Zero physical injury, chemical burns, or vapor inhalation during emergency drill training.", RED_HOT, RGBColor(255, 245, 245)),
        ("90%+", "COST REDUCTION", "Saves ₹50,000–₹1,00,000 per responder per drill by eliminating single-use suit waste.", GREEN_VALID, RGBColor(245, 254, 248)),
        ("1,000+", "DRILLS / DAY", "Infinite virtual repetitions with instantaneous objective grading and replay debriefs.", CARD_BORDER, RGBColor(245, 250, 255))
    ]
    
    kpi_w = Inches(3.75)
    kpi_gap = Inches(0.24)
    
    for idx, (stat, title, desc, col, bg_c) in enumerate(kpis):
        x = Inches(0.8 + idx * (kpi_w.inches + kpi_gap.inches))
        card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.25), kpi_w, Inches(1.65))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_c
        card.line.color.rgb = col
        card.line.width = Pt(1.5)
        
        ctf = card.text_frame
        ctf.margin_left = Inches(0.15)
        ctf.margin_top = Inches(0.1)
        
        cp1 = ctf.paragraphs[0]
        r_stat = cp1.add_run()
        r_stat.text = f"{stat} "
        r_stat.font.size = Pt(20)
        r_stat.font.bold = True
        r_stat.font.color.rgb = col
        
        r_title = cp1.add_run()
        r_title.text = title
        r_title.font.size = Pt(11)
        r_title.font.bold = True
        r_title.font.color.rgb = DARK_NAVY
        
        cp2 = ctf.add_paragraph()
        cp2.text = desc
        cp2.font.size = Pt(8.8)
        cp2.font.color.rgb = TEXT_MUTED
        cp2.space_before = Pt(2)

    # Bottom Left: Target Institutional Reach
    target_card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.05), Inches(6.0), Inches(3.15))
    target_card.fill.solid()
    target_card.fill.fore_color.rgb = SLATE_CARD
    target_card.line.color.rgb = CARD_BORDER
    target_card.line.width = Pt(1.5)
    
    ttf = target_card.text_frame
    ttf.margin_left = Inches(0.2)
    ttf.margin_top = Inches(0.18)
    
    tp = ttf.paragraphs[0]
    tp.text = "🎯 TARGET AUDIENCE & INSTITUTIONAL REACH"
    tp.font.size = Pt(12)
    tp.font.bold = True
    tp.font.color.rgb = DARK_NAVY

    targets = [
        ("National & State Disaster Forces:", "NDRF battalions, SDRF units across all Indian states & Union Territories."),
        ("Petrochemical & Heavy Energy:", "Refineries, offshore rigs & chemical depots (IOCL, ONGC, GAIL, NTPC, NPCIL)."),
        ("Defense & Paramilitary HAZMAT:", "CRPF, CISF airport units, Indian Coast Guard, and Municipal Fire Brigades."),
        ("Environmental Sustainability:", "Zero chemical runoff or hazardous plastic suit waste dumped into landfills.")
    ]

    for title, desc in targets:
        p = ttf.add_paragraph()
        p.space_before = Pt(3)
        r1 = p.add_run()
        r1.text = f"• {title} "
        r1.font.bold = True
        r1.font.size = Pt(9.5)
        r1.font.color.rgb = CARD_BORDER
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(8.8)
        r2.font.color.rgb = TEXT_MUTED

    # Bottom Right: Action Shot with HUD Frame
    action_img = "Bay03_Cinematic_Screenshots/04_FirstPerson_Perspective/04_Trainee_FirstPerson_Containment_Action.png"
    if os.path.exists(action_img):
        s5.shapes.add_picture(action_img, Inches(7.0), Inches(3.05), Inches(5.53), Inches(3.15))
        
        ahud = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.1), Inches(5.85), Inches(5.33), Inches(0.3))
        ahud.fill.solid()
        ahud.fill.fore_color.rgb = DARK_NAVY
        ahud.line.fill.background()
        ahtf = ahud.text_frame
        ahtf.margin_top = Inches(0.04)
        ahp = ahtf.paragraphs[0]
        ahp.text = "⚡ FIRST-PERSON TACTICAL CONTAINMENT ACTION"
        ahp.font.size = Pt(9)
        ahp.font.bold = True
        ahp.font.color.rgb = AMBER_ALERT
        ahp.alignment = PP_ALIGN.CENTER

    # =========================================================
    # SLIDE 6: RESEARCH BASE, STANDARDS & COMPLIANCE
    # =========================================================
    s6 = prs.slides[5]
    for shape in list(s6.shapes):
        if shape.name in ["TextBox 8", "Title 1"]:
            sp = shape._element
            sp.getparent().remove(sp)
        elif shape.has_text_frame and "Your Team Name" in shape.text:
            shape.text = "CBRS-X Team"

    add_header_strip(s6, "RESEARCH BASE, COMPLIANCE & SCIENTIFIC STANDARDS", "AUTHORITATIVE FOUNDATIONS")

    ref_card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.25), Inches(11.73), Inches(4.95))
    ref_card.fill.solid()
    ref_card.fill.fore_color.rgb = SLATE_CARD
    ref_card.line.color.rgb = CARD_BORDER
    ref_card.line.width = Pt(1.5)
    
    rtf = ref_card.text_frame
    rtf.margin_left = Inches(0.25)
    rtf.margin_top = Inches(0.2)
    
    rp = rtf.paragraphs[0]
    rp.text = "📑 OFFICIAL REGULATORY STANDARDS & SCIENTIFIC CITATIONS"
    rp.font.size = Pt(13)
    rp.font.bold = True
    rp.font.color.rgb = DARK_NAVY

    citations = [
        ("1. NDRF Tactical Standard Operating Procedures (SOP)", "Official National Disaster Response Force SOP guidelines for Chemical, Biological, Radiological, and Nuclear (CBRN) First-Responder incident lifecycles (MHA, Govt. of India)."),
        ("2. National Disaster Management Authority (NDMA) Guidelines", "Management of Chemical (Terrorism) Disasters & Industrial Hazardous Material Incident Response Guidelines."),
        ("3. OSHA & NIOSH Hazardous Exposure & PID Standards", "Photoionization Detector (PID) calibration protocols, Permissible Exposure Limits (PEL), and Immediately Dangerous to Life or Health (IDLH) safety caps."),
        ("4. Khronos Group & W3C Web Standards", "OpenXR 1.0 Cross-Platform VR Standard & High-Performance WebGL 2.0 / WebAssembly (WASM) 3D graphics pipeline."),
        ("5. Enterprise Security & Cryptographic Integrity", "Spring Security 6 enterprise authentication, STOMP RFC WebSocket specifications, and SHA-256 tamper-evident digest algorithms.")
    ]

    for title, desc in citations:
        p = rtf.add_paragraph()
        p.space_before = Pt(6)
        r1 = p.add_run()
        r1.text = f"✔ {title}\n"
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = DARK_NAVY
        
        r2 = p.add_run()
        r2.text = f"   {desc}"
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = TEXT_MUTED

    # Ensure max 6 slides
    while len(prs.slides) > 6:
        rId = prs.slides._sldIdLst[6].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[6]

    out_path = "CBRS_X_SIH_2026_Idea_Presentation_Redesigned.pptx"
    prs.save(out_path)
    print(f"Redesigned presentation saved to: {out_path}")

if __name__ == "__main__":
    build_impressive_deck()
