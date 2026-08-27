import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def clean_slide_content(slide, keep_ids):
    for shape in list(slide.shapes):
        if shape.shape_id not in keep_ids and shape.name not in ['Picture 1', 'Picture 10', 'Picture 11', 'Picture 12', 'Oval 8', 'Oval 9', 'Oval 10', 'Oval 11', 'Footer Placeholder 6', 'Slide Number Placeholder 5', 'Rectangle 8', 'Rectangle 9']:
            try:
                sp = shape._element
                sp.getparent().remove(sp)
            except Exception:
                pass

def generate_clean_master_pptx():
    template_path = 'CBRS_X_SIH_2026_Master_Presentation.pptx'
    prs = Presentation(template_path)
    
    DARK_NAVY   = RGBColor(10, 16, 30)
    CARD_BG     = RGBColor(248, 250, 254)
    BORDER_BLUE = RGBColor(0, 120, 215)
    AMBER_ACCENT= RGBColor(230, 95, 0)
    GREEN_ACCENT= RGBColor(16, 160, 100)
    TEXT_MAIN   = RGBColor(15, 23, 42)
    WHITE       = RGBColor(255, 255, 255)

    def set_team_name(slide, name="CBRS-X Team"):
        for shape in slide.shapes:
            if shape.has_text_frame and ("Your Team Name" in shape.text or "Your\nTeam\nName" in shape.text):
                shape.text = name
                if len(shape.text_frame.paragraphs) > 0:
                    shape.text_frame.paragraphs[0].font.size = Pt(10)
                    shape.text_frame.paragraphs[0].font.bold = True
                    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def add_title(slide, title_text):
        for shape in slide.shapes:
            if shape.name == "Title 1" and shape.has_text_frame:
                shape.text = title_text
                shape.text_frame.paragraphs[0].font.size = Pt(22)
                shape.text_frame.paragraphs[0].font.bold = True
                shape.text_frame.paragraphs[0].font.color.rgb = DARK_NAVY
                return
        # If no Title 1, add a clean title box
        tb = slide.shapes.add_textbox(Inches(1.8), Inches(0.4), Inches(8.5), Inches(0.8))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = DARK_NAVY

    # ----------------------------------------------------
    # SLIDE 1: TITLE PAGE
    # ----------------------------------------------------
    s1 = prs.slides[0]
    clean_slide_content(s1, [])
    
    tbox = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(11.73), Inches(1.85))
    tbox.fill.solid()
    tbox.fill.fore_color.rgb = DARK_NAVY
    tbox.line.color.rgb = BORDER_BLUE
    tbox.line.width = Pt(1.5)
    
    ttf = tbox.text_frame
    ttf.margin_left = Inches(0.3)
    ttf.margin_top = Inches(0.18)
    
    tp1 = ttf.paragraphs[0]
    tp1.text = "SMART INDIA HACKATHON 2026"
    tp1.font.size = Pt(13)
    tp1.font.bold = True
    tp1.font.color.rgb = AMBER_ACCENT
    
    tp2 = ttf.add_paragraph()
    tp2.text = "TITLE PAGE // CBRS-X"
    tp2.font.size = Pt(22)
    tp2.font.bold = True
    tp2.font.color.rgb = WHITE
    tp2.space_before = Pt(2)
    
    tp3 = ttf.add_paragraph()
    tp3.text = "Virtual Reality & WebGL CBRN Hazmat Emergency Simulator"
    tp3.font.size = Pt(12)
    tp3.font.color.rgb = RGBColor(190, 215, 245)

    mbox = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.45), Inches(6.0), Inches(4.4))
    mbox.fill.solid()
    mbox.fill.fore_color.rgb = CARD_BG
    mbox.line.color.rgb = BORDER_BLUE
    mbox.line.width = Pt(1.5)
    
    mtf = mbox.text_frame
    mtf.margin_left = Inches(0.25)
    mtf.margin_top = Inches(0.2)
    
    mp1 = mtf.paragraphs[0]
    mp1.text = "📋 PROBLEM STATEMENT DETAILS"
    mp1.font.size = Pt(13)
    mp1.font.bold = True
    mp1.font.color.rgb = DARK_NAVY
    
    s1_items = [
        ("Problem Statement ID:", "SIH260088"),
        ("Problem Statement Title:", "Virtual Reality & WebGL CBRN Emergency Simulator"),
        ("Theme:", "Disaster Management (NDRF / MHA)"),
        ("PS Category:", "Software"),
        ("Institution:", "Kalpataru Institute of Technology (KIT)"),
        ("Team Name:", "CBRS-X Team"),
        ("Team Lead:", "Lohith R C")
    ]
    
    for k, v in s1_items:
        p = mtf.add_paragraph()
        p.space_before = Pt(4)
        r1 = p.add_run()
        r1.text = f"• {k} "
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = BORDER_BLUE
        
        r2 = p.add_run()
        r2.text = f"{v}"
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = TEXT_MAIN

    p1_img = 'scratch/extracted_pdf_images/p1_img0_Image7.png'
    if os.path.exists(p1_img):
        s1.shapes.add_picture(p1_img, Inches(7.1), Inches(2.45), Inches(5.43), Inches(4.4))

    # ----------------------------------------------------
    # SLIDE 2: PROPOSED SOLUTION
    # ----------------------------------------------------
    s2 = prs.slides[1]
    clean_slide_content(s2, [])
    set_team_name(s2)
    add_title(s2, "PROPOSED SOLUTION")

    p2_img = 'scratch/extracted_pdf_images/p2_img2_Image51.png'
    if os.path.exists(p2_img):
        s2.shapes.add_picture(p2_img, Inches(0.8), Inches(1.25), Inches(11.73), Inches(4.95))

    # ----------------------------------------------------
    # SLIDE 3: TECHNICAL APPROACH
    # ----------------------------------------------------
    s3 = prs.slides[2]
    clean_slide_content(s3, [])
    set_team_name(s3)
    add_title(s3, "TECHNICAL APPROACH")

    p3_img = 'scratch/extracted_pdf_images/p3_img2_Image54.png'
    if os.path.exists(p3_img):
        s3.shapes.add_picture(p3_img, Inches(0.8), Inches(1.25), Inches(6.8), Inches(4.95))

    tech_box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.25), Inches(4.73), Inches(4.95))
    tech_box.fill.solid()
    tech_box.fill.fore_color.rgb = CARD_BG
    tech_box.line.color.rgb = BORDER_BLUE
    tech_box.line.width = Pt(1.5)
    
    ttf = tech_box.text_frame
    ttf.margin_left = Inches(0.2)
    ttf.margin_top = Inches(0.18)
    
    tp1 = ttf.paragraphs[0]
    tp1.text = "🛠️ TECHNOLOGY STACK"
    tp1.font.size = Pt(11.5)
    tp1.font.bold = True
    tp1.font.color.rgb = DARK_NAVY
    
    t_stack = [
        ("Unity + WebGL:", "VR simulation environment"),
        ("React + Three.js:", "Trainee monitoring interface"),
        ("Java + Spring Boot:", "Backend & real-time telemetry"),
        ("PostgreSQL:", "Secure storage of trainee data")
    ]
    for k, v in t_stack:
        p = ttf.add_paragraph()
        p.space_before = Pt(1.5)
        r1 = p.add_run()
        r1.text = f"• {k} "
        r1.font.bold = True
        r1.font.size = Pt(9.5)
        r1.font.color.rgb = BORDER_BLUE
        r2 = p.add_run()
        r2.text = v
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = TEXT_MAIN

    tp2 = ttf.add_paragraph()
    tp2.text = "🔄 HOW IT WORKS"
    tp2.font.size = Pt(11.5)
    tp2.font.bold = True
    tp2.font.color.rgb = DARK_NAVY
    tp2.space_before = Pt(5)

    hw = [
        "• Simulate ➔ Capture ➔ Evaluate ➔ Live Monitor ➔ Certify",
        "• Tracks trainee actions in real time",
        "• Automatically calculates a 100-point score",
        "• Live tactical radar view for instructor",
        "• Generates secure, tamper-proof SHA-256 certificate"
    ]
    for line in hw:
        p = ttf.add_paragraph()
        p.space_before = Pt(1.5)
        p.text = line
        p.font.size = Pt(9.2)
        p.font.color.rgb = TEXT_MAIN

    tp3 = ttf.add_paragraph()
    tp3.text = "📊 PERFORMANCE"
    tp3.font.size = Pt(11.5)
    tp3.font.bold = True
    tp3.font.color.rgb = GREEN_ACCENT
    tp3.space_before = Pt(5)

    perf = [
        "✔ 46/46 Tests Passed (JUnit 5 + MockMvc)",
        "✔ Supports up to 10,000 IP requests (LRU Rate Limiter)",
        "✔ Smooth 60 FPS VR/WebGL experience"
    ]
    for line in perf:
        p = ttf.add_paragraph()
        p.space_before = Pt(1.5)
        p.text = line
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = DARK_NAVY

    # ----------------------------------------------------
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # ----------------------------------------------------
    s4 = prs.slides[3]
    clean_slide_content(s4, [])
    set_team_name(s4)
    add_title(s4, "FEASIBILITY AND VIABILITY")

    p4_img = 'scratch/extracted_pdf_images/p4_img1_Image70.png'
    if os.path.exists(p4_img):
        s4.shapes.add_picture(p4_img, Inches(0.8), Inches(1.25), Inches(11.73), Inches(4.95))

    # ----------------------------------------------------
    # SLIDE 5: IMPACT AND BENEFITS
    # ----------------------------------------------------
    s5 = prs.slides[4]
    clean_slide_content(s5, [])
    set_team_name(s5)
    add_title(s5, "IMPACT AND BENEFITS")

    p5_img = 'scratch/extracted_pdf_images/p5_img2_Image73.png'
    if os.path.exists(p5_img):
        s5.shapes.add_picture(p5_img, Inches(0.8), Inches(1.25), Inches(11.73), Inches(4.95))

    # ----------------------------------------------------
    # SLIDE 6: RESEARCH AND REFERENCES
    # ----------------------------------------------------
    s6 = prs.slides[5]
    clean_slide_content(s6, [])
    set_team_name(s6)
    add_title(s6, "RESEARCH AND REFERENCES")

    p6_img = 'scratch/extracted_pdf_images/p6_img2_Image78.png'
    if os.path.exists(p6_img):
        s6.shapes.add_picture(p6_img, Inches(0.8), Inches(1.25), Inches(11.73), Inches(4.95))

    # Remove slide 7 if present
    while len(prs.slides) > 6:
        rId = prs.slides._sldIdLst[6].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[6]

    out_file = "KIT_SPOC_2026_Generated.pptx"
    prs.save(out_file)
    print(f"Pristine PPTX presentation saved successfully: {out_file}")

if __name__ == "__main__":
    generate_clean_master_pptx()
